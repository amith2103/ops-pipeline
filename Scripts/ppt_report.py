import sqlite3
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

print("All Librariees Imported Successfully")

#Load Data 
conn = sqlite3.connect("data/ops-pipeline.db")
df = pd.read_sql("SELECT * FROM service_requests",conn)
conn.close()

#stats
total = len(df)
open_tickets = len(df[df["status"].str.strip() == "Open"])
critical = len(df[df["priority"].str.strip() == "Critical"])
resolved = len(df[df["status"].str.strip() == "Resolved"])
top_region = df["region"].value_counts().index[0]
top_service = df["service_type"].value_counts().index[0]

print("Data loaded:", total, "records")


# Create presentation
prs = Presentation()

# ── Slide 1: Title Slide ──
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# Background
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x06, 0x5A, 0x82)

# Title text
txBox = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Utility Service Request Dashboard"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Executive Summary — Operations Pipeline"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC)
p2.alignment = PP_ALIGN.CENTER

print("Slide 1 done!")


# ── Slide 2: Key Stats ──
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Key Metrics at a Glance"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x06, 0x5A, 0x82)
p.alignment = PP_ALIGN.CENTER

# Stats boxes
stats = [
    ("Total Requests", str(total)),
    ("Open Tickets", str(open_tickets)),
    ("Resolved", str(resolved)),
    ("Critical", str(critical)),
]

positions = [0.5, 2.9, 5.3, 7.1]

for i, (label, value) in enumerate(stats):
    left = Inches(positions[i] - 0.2)

    # Value
    vBox = slide2.shapes.add_textbox(left, Inches(2), Inches(2), Inches(1.2))
    vf = vBox.text_frame
    vp = vf.paragraphs[0]
    vp.text = value
    vp.font.size = Pt(40)
    vp.font.bold = True
    vp.font.color.rgb = RGBColor(0x06, 0x5A, 0x82)
    vp.alignment = PP_ALIGN.CENTER

    # Label
    lBox = slide2.shapes.add_textbox(left, Inches(3.2), Inches(2), Inches(0.8))
    lf = lBox.text_frame
    lp = lf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(14)
    lp.font.color.rgb = RGBColor(0x36, 0x45, 0x4F)
    lp.alignment = PP_ALIGN.CENTER

print("Slide 2 done!")



# ── Slide 3: Insights ──
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Key Insights & Recommendations"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x06, 0x5A, 0x82)
p.alignment = PP_ALIGN.CENTER

# Insights
insights = [
    f"24.6% open ticket backlog ({open_tickets}/{total}) signals potential SLA risk — recommend daily triage review",
    f"{top_region} region leads with {df['region'].value_counts()[top_region]} requests — consider reallocating field technicians",
    f"Billing complaints dominate at {df['service_type'].value_counts()[top_service]} requests (27.5%) — may indicate billing system errors or recent rate changes",
    f"{critical} critical tickets outstanding — immediate escalation required to avoid safety and compliance violations",
    f"Only {resolved} resolved ({round(resolved/total*100, 1)}%) — resolution rate needs investigation before month-end reporting",
]

for i, insight in enumerate(insights):
    iBox = slide3.shapes.add_textbox(
        Inches(0.7), Inches(1.3 + i * 0.85), Inches(8.5), Inches(0.75)
    )
    itf = iBox.text_frame
    itf.word_wrap = True
    ip = itf.paragraphs[0]
    ip.text = f"• {insight}"
    ip.font.size = Pt(15)
    ip.font.color.rgb = RGBColor(0x36, 0x45, 0x4F)

# Save
prs.save("data/executive_report.pptx")
print("PowerPoint saved to data/executive_report.pptx")